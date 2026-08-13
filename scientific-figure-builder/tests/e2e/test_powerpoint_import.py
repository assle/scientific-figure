"""Optional real-PowerPoint end-to-end check.

This test drives the installed Microsoft PowerPoint application through
AppleScript. It verifies that a ``scientific-figure-builder`` SVG exported with
``export_target="ppt"`` can be inserted into PowerPoint, converted from an
imported picture to Office drawing objects, and then ungrouped into many child
shapes.

The test is opt-in because it opens a GUI application and may require the user
to grant PowerPoint access to the working directory once:

    RUN_POWERPOINT_E2E=1 .venv/bin/python -m pytest \
        tests/e2e/test_powerpoint_import.py -q
"""

from __future__ import annotations

import os
import shutil
import subprocess
from pathlib import Path

import matplotlib
import matplotlib.pyplot as plt
import pytest
from pptx import Presentation
from pptx.util import Mm

from figure_tools.export.exporters import export_svg


POWERPOINT_APP = Path("/Applications/Microsoft PowerPoint.app")
POWERPOINT_E2E_DIR = Path(
    os.environ.get(
        "POWERPOINT_E2E_DIR",
        Path.home() / ".cache" / "scientific-figure-powerpoint-e2e",
    )
)


def _powerpoint_available() -> bool:
    return (
        os.environ.get("RUN_POWERPOINT_E2E") == "1"
        and POWERPOINT_APP.exists()
        and shutil.which("osascript") is not None
    )


pytestmark = pytest.mark.skipif(
    not _powerpoint_available(),
    reason=(
        "set RUN_POWERPOINT_E2E=1 and install Microsoft PowerPoint "
        "to run the GUI PowerPoint end-to-end test"
    ),
)


def _run_osascript(script: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        ["osascript"],
        input=script,
        text=True,
        capture_output=True,
        timeout=120,
        check=False,
    )


def _make_svg(path: Path) -> None:
    matplotlib.use("Agg")
    fig, ax = plt.subplots(figsize=(4, 2.5))
    ax.plot([0, 1], [0, 1], linewidth=2)
    ax.set_xlabel("X label")
    ax.set_ylabel("Y label")
    ax.set_title("Title")
    try:
        export_svg(fig, path, export_target="ppt")
    finally:
        plt.close(fig)


def _make_blank_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Mm(180)
    prs.slide_height = Mm(90)
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(path)


def test_powerpoint_imports_and_ungroups_ppt_svg() -> None:
    POWERPOINT_E2E_DIR.mkdir(parents=True, exist_ok=True)
    svg_path = POWERPOINT_E2E_DIR / "figure.svg"
    blank_path = POWERPOINT_E2E_DIR / "blank.pptx"

    _make_svg(svg_path)
    _make_blank_pptx(blank_path)

    svg_text = svg_path.read_text(encoding="utf-8")
    assert "<text" in svg_text

    script = f"""
tell application "Microsoft PowerPoint"
  activate
  open POSIX file "{blank_path}"
  delay 2
  set pres to active presentation
  set sl to slide 1 of pres
  make new picture at sl with properties {{file name:"{svg_path}", left position:0, top:0}}
end tell
tell application "System Events"
  tell process "Microsoft PowerPoint"
    keystroke "a" using command down
    delay 1
    click menu item "取消组合" of menu "排列" of menu bar 1
    delay 1
    if exists button "是(Y)" of front window then
      click button "是(Y)" of front window
      delay 1
    end if
    click menu item "取消组合" of menu "排列" of menu bar 1
  end tell
end tell
tell application "Microsoft PowerPoint"
  set pres to active presentation
  set sl to slide 1 of pres
  set topCount to count of shapes of sl
  set childCount to 0
  repeat with i from 1 to topCount
    try
      set c to count of shapes of shape i of sl
      if c > childCount then set childCount to c
    end try
  end repeat
  set resText to ((topCount as text) & "," & (childCount as text))
  close pres saving no
  return resText
end tell
"""

    proc = _run_osascript(script)
    if proc.returncode != 0:
        pytest.fail(f"PowerPoint AppleScript failed: {proc.stderr.strip()}")

    parts = proc.stdout.strip().split(",")
    assert len(parts) == 2, f"unexpected osascript output: {proc.stdout!r}"
    top_count = int(parts[0])
    child_count = int(parts[1])

    assert top_count >= 2, "SVG was not converted into a group plus a line shape"
    assert child_count >= 10, "ungrouping did not expose the expected child shapes"
