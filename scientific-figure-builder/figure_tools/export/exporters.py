"""Deterministic figure exporters (plan section 15, Phase 2 exit criteria).

Repeated execution must produce byte-identical local artifacts. Timestamps and
other non-deterministic metadata are normalized or stripped.
"""

from __future__ import annotations

import io
import re
from datetime import datetime
from pathlib import Path

import matplotlib

matplotlib.use("Agg")  # non-interactive; safe in headless/CI

from PIL import Image  # noqa: E402

_PDF_METADATA = {
    "Creator": "scientific-figure-builder",
    "Producer": "scientific-figure-builder",
    "CreationDate": datetime(2000, 1, 1),
    "ModDate": datetime(2000, 1, 1),
}

_SVG_HASHSALT = "scientific-figure-builder"
_SVG_DATE_RE = re.compile(rb"<dc:date>.*?</dc:date>", re.DOTALL)
_SVG_COMMENT_RE = re.compile(rb"<!--.*?-->", re.DOTALL)


def export_png(fig, path: Path, dpi: int = 300) -> Path:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi)
    buf.seek(0)
    img = Image.open(buf)
    img.load()
    # Re-encode without metadata for byte-identical output.
    img.save(path, format="PNG", optimize=False, pnginfo=None)
    return path


def export_svg(fig, path: Path) -> Path:
    buf = io.BytesIO()
    with matplotlib.rc_context({"svg.hashsalt": _SVG_HASHSALT}):
        fig.savefig(buf, format="svg")
    data = buf.getvalue()
    data = _SVG_COMMENT_RE.sub(b"", data)
    data = _SVG_DATE_RE.sub(b"", data)
    path.write_bytes(data)
    return path


def export_pdf(fig, path: Path) -> Path:
    fig.savefig(path, format="pdf", metadata=_PDF_METADATA)
    return path


def save_figure(fig, out_dir: Path, basename: str = "plot",
                formats=("png", "svg", "pdf"), dpi: int = 300) -> dict[str, Path]:
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    produced: dict[str, Path] = {}
    if "png" in formats:
        produced["png"] = export_png(fig, out_dir / f"{basename}.png", dpi=dpi)
    if "svg" in formats:
        produced["svg"] = export_svg(fig, out_dir / f"{basename}.svg")
    if "pdf" in formats:
        produced["pdf"] = export_pdf(fig, out_dir / f"{basename}.pdf")
    return produced


def export_pptx(
    placements: list[dict],
    output_path: Path,
    canvas_mm: tuple[float, float],
    title: str | None = None,
) -> Path:
    """Optional editable PPTX (plan sections 2 and 13). Preserves editable text."""
    from pptx import Presentation
    from pptx.util import Mm

    w_mm, h_mm = canvas_mm
    prs = Presentation()
    prs.slide_width = Mm(w_mm)
    prs.slide_height = Mm(h_mm)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    for p in sorted(placements, key=lambda item: item.get("z_order", 0)):
        x, y, bw, bh = p["bbox"]
        slide.shapes.add_picture(
            p["path"],
            Mm(x * w_mm), Mm(y * h_mm),
            Mm(bw * w_mm), Mm(bh * h_mm),
        )

    if title:
        box = slide.shapes.add_textbox(Mm(2), Mm(2), Mm(w_mm - 4), Mm(8))
        box.text_frame.text = title  # editable text

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
